import os
from pptx import Presentation

def generate_presentation(reports_dir, presentation_dir):
    # Ensure output directory exists
    os.makedirs(presentation_dir, exist_ok=True)

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Predictive Maintenance Report"
    slide.placeholders[1].text = "Generated Automatically"

    # Example: Add one slide per report file
    for report_file in os.listdir(reports_dir):
        if report_file.endswith(".txt"):
            filepath = os.path.join(reports_dir, report_file)
            with open(filepath, "r", encoding="utf-8") as f:
                content = f.read()

            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = f"Report: {report_file}"
            slide.placeholders[1].text = content

    # Save as slides.pptx in presentation folder
    output_path = os.path.join(presentation_dir, "slides.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    reports_dir = os.path.join(os.path.dirname(__file__), "..", "reports")
    presentation_dir = os.path.join(os.path.dirname(__file__), "..", "presentation")
    generate_presentation(reports_dir, presentation_dir)
